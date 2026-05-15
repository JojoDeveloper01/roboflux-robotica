from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = Path('presentations/roboflux-presentation-google-slides.pptx')
ROOT_COPY = Path('apresentacao-roboflux.pptx')
PUBLIC_COPY = Path('public/apresentacao.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GRAPHITE = RGBColor(39, 49, 56)
STEEL = RGBColor(95, 110, 116)
CONCRETE = RGBColor(246, 242, 238)
SAGE = RGBColor(184, 195, 152)
COPPER = RGBColor(218, 203, 184)
WHITE = RGBColor(255, 255, 255)
ALLOY = RGBColor(129, 141, 161)


def set_bg(slide, color=CONCRETE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h, size=24, bold=False, color=GRAPHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Aptos'
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def bullet_box(slide, items, x, y, w, h, size=22, color=GRAPHITE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = 'Aptos'
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def pill(slide, text, x, y, w, h, fill, color=GRAPHITE, size=16):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    tf = shp.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    return shp


def title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = SAGE
    add_textbox(slide, 'RoboFlux', 0.85, 0.7, 4, 0.5, 20, True, STEEL)
    add_textbox(slide, title, 0.85, 1.55, 11.4, 1.25, 42, True, GRAPHITE)
    add_textbox(slide, subtitle, 0.9, 3.05, 10.7, 0.7, 21, False, STEEL)
    pill(slide, 'Robotics · Automation · Algorithms', 0.9, 4.3, 4.1, 0.55, COPPER, GRAPHITE)
    add_textbox(slide, 'https://roboflux-robotica.pages.dev/', 0.9, 6.55, 8, 0.35, 13, False, STEEL)
    return slide


def section(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, title, 0.75, 0.45, 11.7, 0.6, 28, True, GRAPHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.78, 1.12, 11.2, 0.45, 15, False, STEEL)
    return slide

# 1
slide = title_slide('Reactive AGV Path Planning', 'A web project that connects structured programming with Robotics and Industrial Automation.')

# 2
slide = section('Project objective', 'The project solves a small but realistic automation problem.')
bullet_box(slide, [
    'Build an interactive web application related to the course area.',
    'Use programming fundamentals: variables, conditionals, loops, functions and DOM manipulation.',
    'Simulate how an AGV calculates a safe route inside a factory floor.',
    'Explain the logic clearly enough to present it, not just show the final page.'
], 0.95, 1.9, 8.4, 3.5, 22)
pill(slide, 'Key idea: code is used to explain a real industrial process.', 0.95, 5.85, 7.2, 0.55, SAGE)

# 3
slide = section('Why this fits Robotics and Industrial Automation', 'Industrial environments change. A robot must react.')
bullet_box(slide, [
    'A factory floor can contain machines, pallets, blocked zones and moving obstacles.',
    'An AGV should not depend on one fixed path only.',
    'The map represents a simplified factory floor using a 10x10 grid.',
    'The algorithm recalculates the route when the map changes.'
], 0.95, 1.85, 7.9, 3.6, 22)
# simple factory diagram
for r in range(4):
    for c in range(6):
        x = 9.0 + c*0.48; y = 1.9 + r*0.48
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.38), Inches(0.38))
        color = SAGE if (r,c)==(3,0) else COPPER if (r,c)==(0,5) else STEEL if (r+c)%4==0 else WHITE
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.color.rgb = CONCRETE
add_textbox(slide, 'Simplified AGV grid', 9.0, 4.2, 3.1, 0.4, 15, True, STEEL)

# 4
slide = section('How the user interacts with the page')
bullet_box(slide, [
    'Open the RoboFlux website.',
    'Look at the AGV grid simulation.',
    'Click “Generate new map” to create a different valid factory layout.',
    'Click “Run algorithm” to watch A* analyse cells and find the route.',
    'Read the explanation sections to understand what is happening.'
], 0.95, 1.6, 8.7, 4.2, 22)
pill(slide, 'The page is interactive, not just informational.', 0.95, 6.05, 5.5, 0.55, COPPER)

# 5
slide = section('The A* algorithm', 'A* finds a path by balancing real cost and estimated distance.')
add_textbox(slide, 'f = g + h', 0.95, 1.85, 4.2, 0.9, 46, True, GRAPHITE)
bullet_box(slide, [
    'g = cost already travelled',
    'h = estimated distance to the destination',
    'f = total priority used to choose the next cell'
], 0.95, 3.05, 6.7, 1.8, 21)
add_textbox(slide, 'Simple explanation', 8.2, 1.85, 3.8, 0.35, 18, True, STEEL)
bullet_box(slide, [
    'The AGV looks at nearby cells.',
    'It avoids blocked cells.',
    'It chooses the most promising option.',
    'It repeats until it reaches the destination.'
], 8.2, 2.35, 4.0, 2.4, 20)

# 6
slide = section('Controlled randomness', 'Random does not mean chaotic.')
bullet_box(slide, [
    'The map generator creates random obstacles.',
    'Before showing the map, the system checks if a route still exists.',
    'If the map is impossible, it tries again.',
    'This makes the project dynamic while keeping the simulation useful.'
], 0.95, 1.85, 8.5, 3.4, 22)
pill(slide, 'Dynamic + valid = better learning experience', 0.95, 5.55, 5.9, 0.55, SAGE)

# 7
slide = section('Technologies used')
techs = [('Astro', 'component structure'), ('HTML', 'page structure'), ('CSS', 'custom visual states'), ('Tailwind', 'layout and design utilities'), ('JavaScript / TypeScript', 'algorithm and DOM interaction'), ('GitHub + Cloudflare Pages', 'version control and deployment')]
for i,(name,desc) in enumerate(techs):
    x = 0.85 + (i%2)*6.1; y = 1.55 + (i//2)*1.45
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.5), Inches(1.0))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE; shp.line.color.rgb = SAGE
    add_textbox(slide, name, x+0.25, y+0.18, 2.4, 0.3, 18, True, GRAPHITE)
    add_textbox(slide, desc, x+0.25, y+0.55, 4.7, 0.3, 13, False, STEEL)

# 8
slide = section('Code organization', 'The algorithm is separated from the interface.')
add_textbox(slide, 'src/lib/pathfinding.ts', 0.95, 1.65, 5.0, 0.45, 22, True, GRAPHITE)
bullet_box(slide, ['Pure pathfinding logic', 'Valid map generation', 'A* route calculation'], 0.95, 2.25, 4.8, 1.7, 19)
add_textbox(slide, 'src/components/AGVPathfinder.astro', 7.05, 1.65, 5.5, 0.45, 22, True, GRAPHITE)
bullet_box(slide, ['Grid rendering', 'Button events', 'Animation and DOM updates'], 7.05, 2.25, 5.0, 1.7, 19)
pill(slide, 'This separation makes the code easier to explain during presentation.', 1.4, 5.55, 10.5, 0.6, COPPER)

# 9
slide = section('What the project demonstrates')
bullet_box(slide, [
    'A useful connection between programming and industrial robotics.',
    'An algorithm with a visible result, not hidden theory.',
    'Good code organization through components and separated logic.',
    'A deployed website that can be accessed by anyone.',
    'A final report available directly from the site.'
], 0.95, 1.55, 9.4, 4.3, 22)

# 10
slide = section('Conclusion and links')
add_textbox(slide, 'RoboFlux turns programming concepts into a practical automation simulation.', 0.95, 1.5, 10.8, 0.7, 28, True, GRAPHITE)
bullet_box(slide, [
    'Website: https://roboflux-robotica.pages.dev/',
    'Report: https://roboflux-robotica.pages.dev/relatorio.pdf',
    'GitHub: https://github.com/JojoDeveloper01/roboflux-robotica'
], 0.95, 3.0, 10.8, 1.8, 20)
pill(slide, 'Call to action: visit the page and run the AGV simulation.', 0.95, 5.65, 8.0, 0.65, SAGE)

# Add slide numbers
for idx, slide in enumerate(prs.slides, start=1):
    add_textbox(slide, f'{idx}/10', 12.25, 6.95, 0.6, 0.25, 10, False, STEEL, PP_ALIGN.RIGHT)

prs.save(OUT)
ROOT_COPY.write_bytes(OUT.read_bytes())
PUBLIC_COPY.write_bytes(OUT.read_bytes())
print(OUT)
